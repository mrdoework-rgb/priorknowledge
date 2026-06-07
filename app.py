import streamlit as st
import google.generativeai as genai
import json
import io
import re
from pptx import Presentation
from pptx.util import Pt

# Constants
PROMPT_TEMPLATE = r"""Act as an expert Science Curriculum Designer. I will provide you with a lesson sequence list below. For each row/lesson provided, I want you to design a set of exactly 8 Prior Knowledge Starter Questions to be used at the very beginning of that lesson.

You must strictly output the response as a single, valid JSON array of objects with no other conversational markdown text before or after the code block.

The JSON structure for each object must follow this format precisely:
{
  "Lesson_title": "The exact title of the lesson provided",
  "sequence_number": The exact number integer provided,
  "recall questions": "A single string containing 4 basic recall questions. Numbered 1. to 4. separated by clean line breaks (\n).",
  "comprehension questions": "A single string containing 3 questions. Numbered 5. to 7. separated by clean line breaks (\n).",
  "application question": "A single string containing exactly 1 question. Numbered as 8."
}

CRITICAL RULES:
1. CUMULATIVE & BRIDGE LOGIC: Questions must ONLY test content from PRIOR lessons.
2. TARGETED SELECTION: Prioritize prerequisite knowledge for the upcoming lesson.
3. ACCESSIBILITY: Use simple, everyday language.
4. BASELINE: For lesson 1, use UK KS2 Science curriculum content.
5. DEDUPLICATION: One JSON object per lesson.

Here is the lesson sequence data:"""

def create_powerpoint(data_list):
    """Generates a PowerPoint and returns it as a BytesIO stream."""
    prs = Presentation()
    for item in data_list:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{item['Lesson_title']} (Lesson {item['sequence_number']})"

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()

        # Combine questions from the stateful text areas
        q_text = f"{item['recall']}\n{item['comp']}\n{item['app']}"
        questions = [q.strip() for q in q_text.split('\n') if q.strip()]

        for i, question in enumerate(questions):
            p = tf.add_paragraph()
            cleaned = re.sub(r'^[\d\.\s]+', '', question)
            p.text = f"{i+1}. {cleaned}"
            p.level = 0
            p.font.size = Pt(24)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# Streamlit UI

st.image("sausages.jpg", use_container_width=True)
st.subheader("The Sausage Factory:- Make a load of Recall Questions for a whole topic where the questions only cover prior knowledge and are  targeted at the content most relevant to the lesson")


# Try to get API key from secrets first, then sidebar
api_key = st.secrets.get("GOOGLE_API_KEY")
sidebar_key = st.sidebar.text_input("Gemini API Key (Override)", type="password")
if sidebar_key:
    api_key = sidebar_key

if api_key:
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-2.5-flash')
st.markdown("Lesson Data in this Format [Put the Learning Objectives into this and copy and paste into the box below](https://priorylewes.sharepoint.com/:x:/s/Science/IQAB3rdK5KCzQJHuqKNH2v9bAV2mWOAtpNOe94_awdPjr08?e=XNcoWG).")
lesson_input = st.text_area("Paste your lesson sequence data here:", height=150)

if "lessons_data" not in st.session_state:
    st.session_state.lessons_data = []

if st.button("Generate Questions"):
    if not api_key:
        st.error("Please provide an API key in the sidebar.")
    elif not lesson_input.strip():
        st.warning("Please enter lesson data first.")
    else:
        with st.spinner("Calling Gemini API..."):
            try:
                full_prompt = PROMPT_TEMPLATE + "\n" + lesson_input
                response = model.generate_content(
                    full_prompt,
                    generation_config=genai.types.GenerationConfig(response_mime_type="application/json")
                )
                parsed = json.loads(response.text)
                # Store in session state for editing
                st.session_state.lessons_data = [
                    {
                        "Lesson_title": l.get("Lesson_title", "N/A"),
                        "sequence_number": l.get("sequence_number", 0),
                        "recall": l.get("recall questions", ""),
                        "comp": l.get("comprehension questions", ""),
                        "app": l.get("application question", "")
                    } for l in parsed
                ]
            except Exception as e:
                st.error(f"API/Parsing Error: {e}")

if st.session_state.lessons_data:
    with st.form("editor_form"):
        st.write("### Edit any of the questions and download the Powerpoint at the bottom")
        for i, lesson in enumerate(st.session_state.lessons_data):
            st.subheader(f"{lesson['Lesson_title']} (Lesson {lesson['sequence_number']})")
            st.session_state.lessons_data[i]['recall'] = st.text_area("Recall Questions", value=lesson['recall'], key=f"rec_{i}")
            st.session_state.lessons_data[i]['comp'] = st.text_area("Comprehension Questions", value=lesson['comp'], key=f"comp_{i}")
            st.session_state.lessons_data[i]['app'] = st.text_area("Application Question", value=lesson['app'], key=f"app_{i}")
        
        submitted = st.form_submit_button("Update and Prepare Download")

    ppt_buffer = create_powerpoint(st.session_state.lessons_data)
    st.download_button(
        label="Download PowerPoint",
        data=ppt_buffer,
        file_name="Lesson_Starter_Questions.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

